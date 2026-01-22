#!/usr/bin/env python3
"""
Generate TOT Presentation from TOT Planning Document
Creates a comprehensive PowerPoint presentation for the 2-day EPS Backend training
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 0, 0)  # EPS Red

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    for p in subtitle_frame.paragraphs:
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)  # EPS Red

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)  # EPS Red

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    # Left title
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)  # EPS Red

    # Left items
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(14)
        p.space_before = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    # Right title
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)  # EPS Red

    # Right items
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(14)
        p.space_before = Pt(4)

    return slide

def add_code_slide(prs, title, code_snippet, language="php"):
    """Add a slide with code example"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)  # EPS Red

    # Code background box
    code_box_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.4), Inches(1.1), Inches(9.2), Inches(5.8)
    )
    code_box_shape.fill.solid()
    code_box_shape.fill.fore_color.rgb = RGBColor(40, 40, 40)  # Dark background
    code_box_shape.line.color.rgb = RGBColor(100, 100, 100)

    # Code text
    code_text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(9), Inches(5.4))
    code_frame = code_text_box.text_frame
    code_frame.word_wrap = True

    p = code_frame.paragraphs[0]
    p.text = code_snippet
    p.font.name = 'Courier New'
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(200, 220, 100)  # Light green for code
    p.line_spacing = 1.0

    return slide

def create_presentation():
    """Create the complete TOT presentation"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs,
        "EPS Backend Web Training",
        "Transfer of Training (TOT)\n2-Day Intensive Course\nJanuary 7-8, 2026"
    )

    # Slide 2: Course Overview
    add_content_slide(prs,
        "Course Overview",
        [
            "• Comprehensive 2-day intensive training",
            "• 16 hours total (8 hours per day)",
            "• Target: Backend developers & system architects",
            "• Framework: Laravel 10 | PHP 8.1+",
            "• 8 hands-on labs with working code",
            "• 50+ real-world code examples",
            "• Complete system: 300+ models, 500+ API routes"
        ]
    )

    # Slide 3: Learning Outcomes - Day 1
    add_content_slide(prs,
        "Day 1 Learning Outcomes",
        [
            "✓ Understand complete project architecture",
            "✓ Work with 300+ Eloquent models",
            "✓ Design complex database relationships",
            "✓ Create and optimize database queries",
            "✓ Build RESTful API endpoints",
            "✓ Transform API responses with Resources",
            "✓ Implement input validation",
            "✓ Optimize queries with eager loading"
        ]
    )

    # Slide 4: Learning Outcomes - Day 2
    add_content_slide(prs,
        "Day 2 Learning Outcomes",
        [
            "✓ Implement JWT & Keycloak SSO authentication",
            "✓ Design permission-based authorization",
            "✓ Create complex business logic services",
            "✓ Handle file uploads & media management",
            "✓ Export data to Excel and PDF",
            "✓ Implement database transactions",
            "✓ Optimize performance & avoid N+1 problems",
            "✓ Create event-driven architecture"
        ]
    )

    # Slide 5: Pre-requisites
    add_two_column_slide(prs,
        "Pre-requisites & Setup",
        "Required Knowledge",
        [
            "• PHP 8.1+ basics",
            "• Laravel 9/10 fundamentals",
            "• Database concepts",
            "• REST API basics",
            "• Object-oriented programming",
            "• Command-line usage"
        ],
        "Required Software",
        [
            "• PHP 8.1 or higher",
            "• Composer",
            "• MySQL/MariaDB 5.7+",
            "• VS Code or IDE",
            "• Git",
            "• Postman/Insomnia"
        ]
    )

    # Slide 6: Day 1 Schedule
    add_content_slide(prs,
        "Day 1: Foundation & Architecture",
        [
            "9:00 - 10:30   Session 1.1: Project Overview & Architecture",
            "10:45 - 12:15  Session 1.2: Database & Eloquent Models (Lab 1.1)",
            "1:00 - 3:00    Session 1.3: Building API Endpoints (Lab 1.2)",
            "3:15 - 4:45    Session 1.4: Query Optimization (Lab 1.3)",
            "4:45 - 5:45    Session 1.5: Wrap-up & Q&A",
            "",
            "Focus: Understanding system, building basic functionality"
        ]
    )

    # Slide 7: Day 2 Schedule
    add_content_slide(prs,
        "Day 2: Advanced Patterns & Implementation",
        [
            "9:00 - 10:30   Session 2.1: Authentication & Authorization (Lab 2.1)",
            "10:45 - 12:15  Session 2.2: Service Layer & Business Logic (Lab 2.2)",
            "1:00 - 2:30    Session 2.3: File Management & Data Export (Lab 2.3)",
            "2:45 - 4:15    Session 2.4: Performance Optimization (Lab 2.4)",
            "4:30 - 5:45    Session 2.5: Advanced Patterns & Best Practices (Lab 2.5)",
            "",
            "Focus: Complex operations, real-world implementation"
        ]
    )

    # Slide 8: Project Architecture
    add_content_slide(prs,
        "Layered Architecture",
        [
            "API Routes (routes/api.php)",
            "         ↓",
            "Controllers (Http/Controllers)",
            "         ↓",
            "Services (App/Services)",
            "         ↓",
            "Models (App/Models)",
            "         ↓",
            "Database (Migrations)"
        ]
    )

    # Slide 9: EPS Modules Overview
    add_content_slide(prs,
        "System Modules",
        [
            "• Course Management (70+ models)",
            "• Exam Management (120+ models)",
            "• Facility Management (30+ models)",
            "• Inspectorate (70+ models)",
            "• User & Agency Management (80+ models)",
            "• System Configuration (40+ models)",
            "• Plus: Payment, Consultation, Digital Safety, Audit"
        ]
    )

    # Slide 10: Eloquent Relationships
    add_content_slide(prs,
        "Core Concepts: Eloquent Relationships",
        [
            "• One-to-Many: Course → Sessions",
            "• Many-to-One (Inverse): Sessions → Course",
            "• Many-to-Many: Through pivot tables",
            "• Polymorphic: Audit logs across models",
            "• Has-Many-Through: Complex chains",
            "• JSON Casting: Auto-encode/decode arrays",
            "• Eager Loading: Avoid N+1 problems"
        ]
    )

    # Slide 10.5: Relationships Code Sample
    add_code_slide(prs,
        "Relationships Code Example",
        """// One-to-Many
public function sessions(): HasMany {
    return $this->hasMany(CourseSession::class);
}

// JSON Casting
protected $casts = [
    'open_to' => 'array',
    'special_to' => 'array',
];

// Eager Loading
$courses = Course::with([
    'subCategory',
    'createdBy',
    'sessions'
])->paginate(15);"""
    )

    # Slide 11: API Design
    add_content_slide(prs,
        "RESTful API Design",
        [
            "• Resource-oriented design",
            "• HTTP methods: GET, POST, PUT, DELETE",
            "• Proper HTTP status codes",
            "• Request/response consistency",
            "• Pagination for large datasets",
            "• API resources for transformation",
            "• Form request validation",
            "• 500+ endpoints organized by module"
        ]
    )

    # Slide 11.5: API Controller Code Sample
    add_code_slide(prs,
        "Controller & Validation Example",
        """class CourseCategoryController extends Controller {
    public function store(StoreCategoryRequest $request) {
        $category = CourseCategory::create([
            ...$request->validated(),
            'created_by' => auth()->id(),
        ]);
        return new CourseCategoryResource($category);
    }

    public function index() {
        $categories = CourseCategory::with('createdBy')
            ->when(request('search'), fn($q) =>
                $q->where('name', 'like', '%'.search().'%'))
            ->paginate(15);
        return CourseCategoryResource::collection($categories);
    }
}"""
    )

    # Slide 12: Lab 1.1
    add_content_slide(prs,
        "Lab 1.1: Model Creation & Relationships",
        [
            "Create CoursePrerequisite model with:",
            "• Database migration (foreign keys, constraints)",
            "• Bi-directional relationships",
            "• JSON casting for array storage",
            "• Auditable interface for change tracking",
            "• Factory for testing",
            "",
            "Duration: 30 minutes"
        ]
    )

    # Slide 13: Lab 1.2 - Overview
    add_content_slide(prs,
        "Lab 1.2: Complete API Endpoint",
        [
            "Build CourseCategory CRUD endpoint including:",
            "• Model with migrations",
            "• Form request validation",
            "• API resource transformation",
            "• Complete controller with CRUD operations",
            "• Route registration",
            "• Permission checking",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 13.5: Lab 1.2 - Code Sample (Model)
    add_code_slide(prs,
        "Lab 1.2: Model Example",
        """class CourseCategory extends Model {
    use SoftDeletes, HasFactory;
    use \\OwenIt\\Auditing\\Auditable;

    protected $fillable = [
        'name', 'description', 'status', 'created_by'
    ];

    public function createdBy(): BelongsTo {
        return $this->belongsTo(User::class);
    }
}"""
    )

    # Slide 14: Lab 1.3
    add_content_slide(prs,
        "Lab 1.3: Query Optimization",
        [
            "Convert inefficient queries to optimized versions:",
            "• Identify N+1 query problems",
            "• Implement eager loading with with()",
            "• Use withCount() for aggregates",
            "• Apply query constraints",
            "• Benchmark performance improvements",
            "• Reduce 300+ queries to 3-5 queries",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 14.5: Query Optimization Code Sample
    add_code_slide(prs,
        "Lab 1.3: Query Optimization Example",
        """// INEFFICIENT - N+1 Problem (300+ queries)
$courses = Course::all();
foreach ($courses as $course) {
    echo $course->subCategory->name;  // +100 queries
}

// EFFICIENT - Eager Loading (3 queries)
$courses = Course::with([
    'subCategory',
    'createdBy',
    'sessions',
    'courseCalendars.participants'
])->paginate(20);

// With aggregates
$courses = Course::withCount(['sessions'])
    ->where('status', 'active')
    ->get();"""
    )

    # Slide 15: Authentication Deep Dive
    add_content_slide(prs,
        "Authentication & Security",
        [
            "JWT Authentication:",
            "  • Token generation and validation",
            "  • Token refresh strategy",
            "  • Expiration handling",
            "",
            "Keycloak SSO Integration:",
            "  • Single sign-on configuration",
            "  • User synchronization",
            "  • Multi-system authentication"
        ]
    )

    # Slide 15.5: JWT Authentication Code Sample
    add_code_slide(prs,
        "JWT Authentication Example",
        """public function login(Request $request) {
    $credentials = $request->validate([
        'email' => 'required|email',
        'password' => 'required',
    ]);

    $token = auth('api')->attempt($credentials);

    if (!$token) {
        return response()->json(
            ['message' => 'Invalid credentials'], 401
        );
    }

    return response()->json([
        'access_token' => $token,
        'token_type' => 'Bearer',
        'expires_in' => auth()->factory()->getTTL() * 60,
    ]);
}"""
    )

    # Slide 16: Authorization & Permissions
    add_content_slide(prs,
        "Authorization with Spatie Permission",
        [
            "Role-Based Access Control (RBAC):",
            "  • Define roles (Admin, Manager, User)",
            "  • Assign permissions to roles",
            "  • Check permissions in controllers",
            "",
            "Implementation:",
            "  • Middleware-based authorization",
            "  • Policy-based authorization",
            "  • Custom authorization logic"
        ]
    )

    # Slide 16.5: Authorization Code Sample
    add_code_slide(prs,
        "RBAC Implementation Example",
        """// Assign role to user
$user->assignRole('course-manager');

// Assign permission to role
$role->givePermissionTo('create-course');

// Check permission in controller
if (auth()->user()->can('create-course')) {
    // Create course
}

// Using middleware
Route::middleware('permission:edit-course')
    ->post('/courses/{id}', [...]);

// Using policy
class CoursePolicy {
    public function edit(User $user, Course $course) {
        return $user->id === $course->created_by
            || $user->isAdmin();
    }
}"""
    )

    # Slide 17: Lab 2.1
    add_content_slide(prs,
        "Lab 2.1: Role-Based Access Control",
        [
            "Implement permission-protected endpoint:",
            "• Define permissions (approve-course, reject-course)",
            "• Assign to roles",
            "• Create authorization policies",
            "• Protect endpoints with middleware",
            "• Test with and without permissions",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 18: Service Layer Pattern
    add_content_slide(prs,
        "Service Layer Pattern",
        [
            "Separation of Concerns:",
            "  • Business logic in services, not controllers",
            "  • Reusable across multiple endpoints",
            "  • Easier testing and maintenance",
            "",
            "Multi-step Operations:",
            "  • Validation before execution",
            "  • Database transactions (atomic)",
            "  • Side effects (notifications, cache)",
            "  • Comprehensive error handling"
        ]
    )

    # Slide 18.5: Service Layer Code Sample
    add_code_slide(prs,
        "Service Layer Example",
        """class CourseParticipantService {
    public function registerParticipant($userId, $calId) {
        return DB::transaction(function () use
            ($userId, $calId) {

            // Step 1: Validate eligibility
            $this->validateEligibility($userId, $calId);

            // Step 2: Create participant
            $participant = CourseCalendarParticipant::create([
                'user_id' => $userId,
                'course_calendar_id' => $calId,
            ]);

            // Step 3: Initialize attendance
            $this->initializeAttendance($participant);

            // Step 4: Send notification
            $this->notifyParticipant($participant);

            // Step 5: Clear cache
            cache()->forget('participants_'.$calId);

            return $participant;
        });
    }
}"""
    )

    # Slide 19: Lab 2.2
    add_content_slide(prs,
        "Lab 2.2: Complex Business Service",
        [
            "Course Completion Service with:",
            "• Multi-step process (validation → processing → completion)",
            "• Database transactions for atomicity",
            "• Grade calculation",
            "• Certificate generation",
            "• Notification sending",
            "• Cache invalidation",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 20: File Management
    add_content_slide(prs,
        "File Management & Uploads",
        [
            "Spatie Media Library:",
            "  • File upload handling",
            "  • Media collections",
            "  • File validation (MIME type, size)",
            "  • URL generation",
            "",
            "Usage:",
            "  • Course thumbnails",
            "  • Course materials (PDFs, documents)",
            "  • Certificates and reports"
        ]
    )

    # Slide 21: Data Export
    add_content_slide(prs,
        "Excel & PDF Export",
        [
            "Excel Export with Maatwebsite:",
            "  • FromQuery interface for large datasets",
            "  • Data mapping and transformation",
            "  • Custom formatting and styling",
            "",
            "PDF Generation:",
            "  • DomPDF integration",
            "  • Template rendering",
            "  • Watermarking and signatures",
            "  • Memory-efficient streaming"
        ]
    )

    # Slide 22: Lab 2.3
    add_content_slide(prs,
        "Lab 2.3: Excel Export",
        [
            "Create Excel export with:",
            "• Query optimization for large datasets",
            "• Data mapping and calculations",
            "• Custom column formatting",
            "• Headings and auto-sizing",
            "• Real-time metric calculations",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 23: Performance Optimization
    add_content_slide(prs,
        "Performance Optimization",
        [
            "Query Optimization:",
            "  • Eager loading to avoid N+1",
            "  • Query analysis and debugging",
            "  • Database indexes",
            "",
            "Caching Strategies:",
            "  • Cache tags for granular control",
            "  • Smart cache invalidation",
            "  • Redis for performance"
        ]
    )

    # Slide 23.5: Caching Code Sample
    add_code_slide(prs,
        "Caching Strategy Example",
        """// Cache with remember
public function getCategories() {
    return cache()->remember(
        'course_categories_all',
        now()->addHours(24),
        fn() => CourseCategory::with('createdBy')->get()
    );
}

// Invalidate on update
public function updateCategory($category, Request $req) {
    $category->update($req->validated());
    cache()->forget('course_categories_all');
    cache()->tags(['categories'])->flush();

    return new CategoryResource($category);
}

// Tags for granular control
cache()->tags(['course', "cat-{$catId}"])
    ->remember("courses_{$catId}", 12*60,
        fn() => Course::where(...)->get());"""
    )

    # Slide 24: Lab 2.4
    add_content_slide(prs,
        "Lab 2.4: Caching Strategy",
        [
            "Implement caching for performance:",
            "• Cache frequently accessed data",
            "• Use cache tags for organization",
            "• Invalidate on data changes",
            "• Monitor cache effectiveness",
            "• Test performance improvements",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 25: Advanced Patterns
    add_content_slide(prs,
        "Advanced Patterns",
        [
            "Polymorphic Relationships:",
            "  • Single table for multiple model types",
            "  • Audit logging across all models",
            "",
            "Event-Driven Architecture:",
            "  • Decouple side effects from models",
            "  • Event listeners for notifications",
            "  • Async job processing"
        ]
    )

    # Slide 25.5: Observer Pattern Code Sample
    add_code_slide(prs,
        "Observer Pattern Example",
        """// Create observer
php artisan make:observer CourseObserver --model=Course

// Implement observer
class CourseObserver {
    public function created(Course $course): void {
        Mail::to(config('mail.admin'))
            ->send(new CourseCreatedMail($course));
        activity()->log('Course created');
    }

    public function updated(Course $course): void {
        activity()->withProperties($course->getChanges())
            ->log('Course updated');
        cache()->forget('course_'.$course->id);
    }
}

// Register in AppServiceProvider
Course::observe(CourseObserver::class);"""
    )

    # Slide 26: Lab 2.5
    add_content_slide(prs,
        "Lab 2.5: Observer Pattern",
        [
            "Implement model observer for:",
            "• Automatic email notifications",
            "• Activity logging",
            "• Cache invalidation",
            "• Related data cleanup",
            "",
            "Decouples side effects from business logic",
            "",
            "Duration: 45 minutes"
        ]
    )

    # Slide 27: Best Practices
    add_content_slide(prs,
        "Best Practices Summary",
        [
            "Code Organization:",
            "  • Thin controllers, fat services",
            "  • Business logic separated from models",
            "",
            "Database:",
            "  • Always use migrations",
            "  • Add indexes strategically",
            "  • Use soft deletes for audit trail",
            "",
            "Performance:",
            "  • Eager load relationships",
            "  • Implement caching",
            "  • Optimize queries"
        ]
    )

    # Slide 28: Security Best Practices
    add_content_slide(prs,
        "Security Best Practices",
        [
            "• Always validate input (Form Requests)",
            "• Use policy classes for authorization",
            "• Hash passwords - never store plaintext",
            "• Escape output to prevent XSS",
            "• Use HTTPS in production",
            "• Implement rate limiting",
            "• Sanitize user input",
            "• Use CSRF tokens for web routes"
        ]
    )

    # Slide 29: Testing Strategy
    add_content_slide(prs,
        "Testing Approach",
        [
            "Unit Tests:",
            "  • Test model relationships and scopes",
            "  • Test business logic methods",
            "",
            "Feature Tests:",
            "  • Test API endpoints",
            "  • Test authorization and permissions",
            "  • Test error handling",
            "",
            "Use factories for test data",
            "Use RefreshDatabase trait for isolation"
        ]
    )

    # Slide 30: Final Project Assignment
    add_content_slide(prs,
        "Final Project: CourseApproval Module",
        [
            "Build complete module including:",
            "• Database design with relationships",
            "• Complete CRUD API endpoints",
            "• Permission-based access control",
            "• Multi-step approval workflow",
            "• Email notifications",
            "• Activity logging",
            "",
            "Expected time: 3-4 hours | Points: 100"
        ]
    )

    # Slide 31: Evaluation Criteria
    add_content_slide(prs,
        "Project Evaluation",
        [
            "Code Quality & Standards       20 points",
            "Functionality & Completeness   25 points",
            "Documentation                  15 points",
            "Performance Optimization       15 points",
            "Security Implementation        15 points",
            "Testing Coverage               10 points",
            "",
            "Total: 100 points"
        ]
    )

    # Slide 32: Course Statistics
    add_content_slide(prs,
        "Course Statistics",
        [
            "Total Hours: 16 (2 days × 8 hours)",
            "Practical Labs: 8",
            "Code Examples: 50+",
            "Models Covered: 15+",
            "Key Patterns: 12",
            "Workshops: 4",
            "Assessment Tasks: 6",
            "Success Rate Target: 80%+"
        ]
    )

    # Slide 33: Tools & Resources
    add_content_slide(prs,
        "Tools & Resources",
        [
            "Development:",
            "  • VS Code | Postman | Tinker | Git",
            "",
            "Documentation:",
            "  • Laravel 10 official docs",
            "  • Spatie packages documentation",
            "  • API collection for Postman",
            "",
            "References:",
            "  • HANDS_ON_TRAINING_GUIDE.md",
            "  • Sample project code"
        ]
    )

    # Slide 34: Post-Course Follow-up
    add_content_slide(prs,
        "After Training",
        [
            "Week 1: Review labs, start final project",
            "Week 2-3: Complete project assignment",
            "Week 3-4: Code reviews & knowledge sharing",
            "",
            "Ongoing:",
            "  • Mentor new team members",
            "  • Contribute to project improvements",
            "  • Document advanced patterns",
            "  • Share best practices"
        ]
    )

    # Slide 35: Q&A Slide
    add_title_slide(prs,
        "Questions & Discussion",
        "Ready to start your Transfer of Training journey!"
    )

    # Save presentation
    output_file = r"c:\Users\User\Documents\laragon\www\eps-be-web\EPS_TOT_Training_2Days.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Installing python-pptx...")
        import subprocess
        subprocess.check_call(["pip", "install", "python-pptx"])
        create_presentation()
    except Exception as e:
        print(f"Error: {e}")
