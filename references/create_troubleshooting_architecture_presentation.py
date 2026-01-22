#!/usr/bin/env python3
"""
EPS Backend Web - Troubleshooting & Architecture Presentation Generator
Generates a comprehensive PowerPoint presentation covering:
- Troubleshooting Guide (20 categories, 60+ issues)
- Project Architecture (15 sections, complete technical reference)

Color Scheme: EPS Red (204, 0, 0) and White
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# EPS Brand Colors
EPS_RED = RGBColor(204, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(242, 242, 242)

def create_presentation():
    """Create the presentation object"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle):
    """Add title slide with red background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = EPS_RED

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = f"January 7, 2026"
    date_p.font.size = Pt(14)
    date_p.font.color.rgb = LIGHT_GRAY
    date_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add content slide with red title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = EPS_RED

    # Title underline
    slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9), Inches(0.02)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = EPS_RED
    slide.shapes[-1].line.color.rgb = EPS_RED

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

def add_two_column_slide(prs, title, left_items, right_items):
    """Add two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = EPS_RED

    # Title underline
    slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9), Inches(0.02)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = EPS_RED
    slide.shapes[-1].line.color.rgb = EPS_RED

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

def generate_presentation():
    """Generate the complete presentation"""
    prs = create_presentation()

    # ============ TITLE SLIDE ============
    add_title_slide(prs,
        "EPS Backend Web",
        "Troubleshooting & Architecture Guide")

    # ============ TABLE OF CONTENTS ============
    add_content_slide(prs, "Presentation Overview", [
        "• Part 1: Project Architecture (15 sections)",
        "• Part 2: Troubleshooting Guide (20 categories, 60+ issues)",
        "• Quick Reference & Solutions",
        "• Best Practices & Optimization",
        "• Security & Deployment"
    ])

    # ============ PART 1: PROJECT ARCHITECTURE ============

    add_title_slide(prs,
        "Part 1: Project Architecture",
        "Complete Technical Reference")

    add_content_slide(prs, "System Overview", [
        "• Framework: Laravel 10 REST API",
        "• 300+ Eloquent Models",
        "• 500+ API Routes",
        "• 6 Major Business Modules",
        "• Multi-agency Support",
        "• JWT + Keycloak SSO Authentication"
    ])

    add_content_slide(prs, "Technology Stack", [
        "• Backend: Laravel 10.x + PHP 8.1+",
        "• Database: MySQL 8.0 / MariaDB 10.x",
        "• Cache: Redis 6.x",
        "• Authentication: JWT (tymon/jwt-auth v2.1)",
        "• SSO: Keycloak (robsontenorio/laravel-keycloak-guard)",
        "• Packages: Spatie (Permission, Auditing, Media Library)"
    ])

    add_two_column_slide(prs, "Architecture Layers",
        [
            "Presentation Layer:",
            "• API Controllers",
            "• Request Validation",
            "• Response Formatting",
            "",
            "Business Logic Layer:",
            "• Service Classes",
            "• Transaction Management",
            "• Workflow Orchestration"
        ],
        [
            "Data Access Layer:",
            "• Eloquent Models",
            "• Query Scopes",
            "• Relationships",
            "",
            "Infrastructure:",
            "• Helpers & Traits",
            "• Observers & Events",
            "• Jobs & Notifications"
        ]
    )

    add_content_slide(prs, "Directory Structure", [
        "• app/Controllers/API/ - API endpoints",
        "• app/Models/ - Database models (300+)",
        "• app/Services/ - Business logic",
        "• app/Jobs/ - Queue jobs",
        "• app/Mail/ - Email templates",
        "• database/migrations/ - Schema changes",
        "• routes/api.php - 500+ API routes"
    ])

    add_two_column_slide(prs, "Database Architecture",
        [
            "Core Tables:",
            "• Users & Agencies",
            "• Courses & Sessions",
            "• Course Participants",
            "• Exams & Questions",
            "• Facilities & Bookings",
            "• Audit & Media"
        ],
        [
            "Design Principles:",
            "• Normalization (3NF)",
            "• Strategic Indexing",
            "• Soft Deletes",
            "• Timestamps",
            "• UUID Support",
            "• Full ACID Compliance"
        ]
    )

    add_content_slide(prs, "Authentication & Authorization", [
        "JWT Authentication:",
        "• Token-based, stateless",
        "• HS256 signature algorithm",
        "• Secure secret configuration",
        "",
        "Keycloak SSO:",
        "• OAuth 2.0 / OpenID Connect",
        "• Centralized user management",
        "• Role synchronization"
    ])

    add_content_slide(prs, "RBAC (Role-Based Access Control)", [
        "Roles:",
        "• super_admin, course_manager, course_instructor",
        "• course_student, facility_manager, auditor",
        "",
        "Authorization Methods:",
        "• Middleware: middleware('permission:edit-courses')",
        "• Policies: authorize('update', $course)",
        "• Gates: gate('edit-course')"
    ])

    add_content_slide(prs, "API Architecture", [
        "RESTful Design:",
        "• Standard resource operations (CRUD)",
        "• Nested resources support",
        "• Consistent response format",
        "",
        "Features:",
        "• Pagination & filtering",
        "• API versioning support",
        "• Rate limiting (60 requests/min)"
    ])

    add_two_column_slide(prs, "Business Modules",
        [
            "1. Course Management:",
            "• Course CRUD",
            "• Session scheduling",
            "• Participant enrollment",
            "• Evaluation & feedback",
            "",
            "2. Examination System:",
            "• Exam creation",
            "• Question bank",
            "• Timed exams & grading"
        ],
        [
            "3. Facility Management:",
            "• Venue booking",
            "• Resource allocation",
            "• Maintenance tracking",
            "",
            "4. Inspectorate Module:",
            "• Audit scheduling",
            "• Compliance tracking",
            "• Finding recording"
        ]
    )

    add_content_slide(prs, "Design Patterns", [
        "• Repository Pattern - Data access abstraction",
        "• Service Layer - Business logic encapsulation",
        "• Observer Pattern - Model event handling",
        "• Factory Pattern - Model creation",
        "• Strategy Pattern - Notification channels",
        "• Dependency Injection - Loose coupling"
    ])

    add_content_slide(prs, "Caching Strategy", [
        "Multi-Layer Caching:",
        "• Application Cache (Redis) - Query results",
        "• Session Cache - User sessions",
        "• Route/Config Cache - Production optimization",
        "• OPcache - PHP opcode caching",
        "",
        "Cache Invalidation:",
        "• Tag-based flushing",
        "• Remember keys",
        "• Manual clearing"
    ])

    add_content_slide(prs, "Queue & Job Processing", [
        "Job System:",
        "• Queued in Redis",
        "• Processed by workers",
        "• Retry up to 3 times",
        "",
        "Job Types:",
        "• SendEmailJob, ProcessReportJob",
        "• ExportDataJob, ImportDataJob",
        "• CleanupFilesJob"
    ])

    add_two_column_slide(prs, "File Storage Architecture",
        [
            "Storage Disks:",
            "• Local - Development",
            "• Public - Web accessible",
            "• S3 - AWS cloud storage",
            "",
            "Media Collections:",
            "• Thumbnails",
            "• Courses materials",
            "• Exam attachments"
        ],
        [
            "Upload Flow:",
            "• Validate file",
            "• Scan for viruses",
            "• Generate unique filename",
            "• Store to disk",
            "• Create media record",
            "• Generate thumbnail"
        ]
    )

    add_content_slide(prs, "Security Architecture", [
        "8 Security Layers:",
        "1. Authentication - JWT & Keycloak",
        "2. Authorization - RBAC & Policies",
        "3. Input Validation - Form requests",
        "4. SQL Injection Prevention - Eloquent ORM",
        "5. XSS Prevention - Blade escaping",
        "6. CSRF Protection - Automatic tokens",
        "7. Rate Limiting - Throttle middleware",
        "8. Sensitive Data - Encryption & hiding"
    ])

    add_content_slide(prs, "Deployment Architecture", [
        "Development:",
        "• Local (Laragon/XAMPP)",
        "• PHP, MySQL, Redis, Keycloak (Docker)",
        "",
        "Production:",
        "• Load Balancer",
        "• Multiple app servers (Nginx + PHP-FPM)",
        "• Database cluster (Master-Replica)",
        "• Redis cluster",
        "• Keycloak HA setup"
    ])

    # ============ PART 2: TROUBLESHOOTING GUIDE ============

    add_title_slide(prs,
        "Part 2: Troubleshooting Guide",
        "60+ Issues Across 20 Categories")

    add_content_slide(prs, "Installation & Setup Issues", [
        "1. Composer Install Fails",
        "   • Solution: Increase memory limit, clear cache",
        "",
        "2. Application Key Not Set",
        "   • Solution: php artisan key:generate",
        "",
        "3. Storage Link Not Created",
        "   • Solution: php artisan storage:link",
        "",
        "4. Permission Denied Errors",
        "   • Solution: icacls storage /grant Users:F /T"
    ])

    add_content_slide(prs, "Database Problems", [
        "1. Connection Refused",
        "   • Solution: Check MySQL running, verify .env settings",
        "",
        "2. Database Does Not Exist",
        "   • Solution: CREATE DATABASE eps_be_web",
        "",
        "3. Migration Fails",
        "   • Solution: Check migration order, drop conflicting tables",
        "",
        "4. Foreign Key Constraint Fails",
        "   • Solution: Disable FK checks, ensure parent records exist"
    ])

    add_content_slide(prs, "Authentication Issues", [
        "1. JWT Token Not Generated",
        "   • Solution: php artisan jwt:secret",
        "",
        "2. Token Expired / Invalid",
        "   • Solution: Implement token refresh, increase TTL",
        "",
        "3. User Not Authenticated",
        "   • Solution: Verify token format, check middleware",
        "",
        "4. Password Not Matching",
        "   • Solution: Use Hash::check(), verify hashing method"
    ])

    add_content_slide(prs, "Keycloak SSO Issues", [
        "1. Cannot Connect to Keycloak",
        "   • Solution: docker ps, docker logs, verify URL",
        "",
        "2. Invalid Client / Client Not Found",
        "   • Solution: Verify client exists, regenerate secret",
        "",
        "3. Invalid Token / Token Validation Failed",
        "   • Solution: Update public key, clear cache",
        "",
        "4. CORS Errors",
        "   • Solution: Add web origins in Realm Settings"
    ])

    add_content_slide(prs, "API & Route Issues", [
        "1. 404 Not Found",
        "   • Solution: php artisan route:list, verify route exists",
        "",
        "2. 405 Method Not Allowed",
        "   • Solution: Check HTTP method matches route",
        "",
        "3. Route Model Binding Not Working",
        "   • Solution: Verify parameter name matches model",
        "",
        "4. Validation Errors Not Returned",
        "   • Solution: Check Accept header, implement error handler"
    ])

    add_content_slide(prs, "File Upload & Media Issues", [
        "1. File Upload Fails",
        "   • Solution: Increase upload limits in php.ini",
        "",
        "2. Spatie Media Library Not Working",
        "   • Solution: Ensure model has HasMedia, run migrations",
        "",
        "3. Uploaded Files Not Accessible",
        "   • Solution: Create storage link, check permissions",
        "",
        "4. S3 Upload Fails",
        "   • Solution: Verify AWS credentials, bucket permissions"
    ])

    add_content_slide(prs, "Performance Issues", [
        "1. N+1 Query Problem",
        "   • Solution: Use eager loading with('relation')",
        "",
        "2. Memory Exhausted",
        "   • Solution: Increase memory limit, use chunking",
        "",
        "3. Slow Queries",
        "   • Solution: Add indexes, analyze queries, use caching",
        "",
        "4. Too Many Database Connections",
        "   • Solution: Check active connections, increase max_connections"
    ])

    add_content_slide(prs, "Cache & Session Problems", [
        "1. Cache Not Working",
        "   • Solution: Check CACHE_DRIVER, ensure Redis running",
        "",
        "2. Session Not Persisting",
        "   • Solution: Check SESSION_DRIVER, run migrations",
        "",
        "3. Redis Connection Failed",
        "   • Solution: Start Redis, verify port 6379",
        "",
        "All: Use php artisan cache:clear, config:clear"
    ])

    add_content_slide(prs, "Email & Notification Issues", [
        "1. Emails Not Sending",
        "   • Solution: Check MAIL_MAILER config, use app password",
        "",
        "2. Queue Not Processing",
        "   • Solution: Start queue worker, check failed jobs",
        "",
        "Queue Commands:",
        "• php artisan queue:work",
        "• php artisan queue:failed",
        "• php artisan queue:retry all"
    ])

    add_content_slide(prs, "Permission & Role Issues", [
        "1. Permission Denied / 403 Forbidden",
        "   • Solution: Assign role/permission, verify middleware",
        "",
        "2. Roles Not Syncing",
        "   • Solution: Cache reset, verify tables, check User trait",
        "",
        "Commands:",
        "• php artisan permission:cache-reset",
        "• php artisan cache:clear"
    ])

    add_content_slide(prs, "Troubleshooting Checklist", [
        "When something goes wrong, follow these steps:",
        "",
        "1. php artisan cache:clear",
        "2. tail -f storage/logs/laravel.log",
        "3. php artisan about (check environment)",
        "4. composer dump-autoload",
        "5. php artisan migrate (check status)",
        "6. Set permissions: icacls storage /grant Users:F /T",
        "7. Restart services (web server, queue, Redis)"
    ])

    add_content_slide(prs, "Debugging Tools & Resources", [
        "Built-in Tools:",
        "• Laravel Telescope - request debugging",
        "• Laravel Debugbar - bar debugging",
        "• Laravel Log Viewer - log inspection",
        "",
        "External Resources:",
        "• Laravel Documentation: laravel.com/docs",
        "• Stack Overflow: [laravel] tag",
        "• API Testing: Postman, Insomnia"
    ])

    add_content_slide(prs, "Quick Reference: Common Errors", [
        "SQLSTATE[HY000] [2002] - DB connection refused",
        "SQLSTATE[42S01] - Table already exists",
        "SQLSTATE[23000] - Foreign key constraint fails",
        "401 Unauthorized - Invalid/missing JWT token",
        "403 Forbidden - Insufficient permissions",
        "404 Not Found - Route/resource not found",
        "422 Unprocessable - Validation failed",
        "500 Server Error - Check laravel.log"
    ])

    add_content_slide(prs, "Best Practices Summary", [
        "Development:",
        "• Use eager loading to prevent N+1 queries",
        "• Implement caching for expensive operations",
        "• Validate all user inputs",
        "• Write meaningful error messages",
        "",
        "Production:",
        "• Cache configuration & routes",
        "• Monitor queue workers",
        "• Regular database backups",
        "• Monitor application logs"
    ])

    # ============ CONCLUSION ============

    add_title_slide(prs,
        "Questions?",
        "Reference: TROUBLESHOOTING_GUIDE.md & PROJECT_ARCHITECTURE.md")

    return prs

def main():
    """Main execution"""
    print("Generating EPS Troubleshooting & Architecture Presentation...")
    prs = generate_presentation()

    filename = f"EPS_Troubleshooting_Architecture_{datetime.now().strftime('%Y%m%d')}.pptx"
    prs.save(filename)

    print(f"✓ Presentation created: {filename}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Color scheme: EPS Red (204, 0, 0) and White")

if __name__ == "__main__":
    main()
